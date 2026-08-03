from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
import copy

# ── Colour Palette (SAP-inspired) ──────────────────────────────────────────
SAP_BLUE      = RGBColor(0x00, 0x6D, 0xD7)   # primary blue
SAP_DARK      = RGBColor(0x00, 0x33, 0x66)   # dark navy
SAP_TEAL      = RGBColor(0x00, 0x79, 0x8C)
SAP_GREEN     = RGBColor(0x10, 0x7E, 0x3E)
SAP_AMBER     = RGBColor(0xE9, 0x73, 0x00)
SAP_RED       = RGBColor(0xBB, 0x00, 0x21)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY    = RGBColor(0xF4, 0xF6, 0xF8)
MID_GREY      = RGBColor(0xCB, 0xD6, 0xE2)
TEXT_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_GOLD   = RGBColor(0xF0, 0xAB, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.fill.solid()
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=12, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=11, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
             space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.25, fill=SAP_DARK)
    add_rect(slide, 0, 1.25, 13.33, 0.06, fill=SAP_BLUE)
    add_text(slide, title, 0.4, 0.12, 12.0, 0.65,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.72, 12.0, 0.45,
                 size=14, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)

def footer_bar(slide, note="Hydrocarbon Tank Stock Reconciliation  |  SAP BTP + IS-OIL HPM"):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill=SAP_DARK)
    add_text(slide, note, 0.3, 7.22, 12.7, 0.25,
             size=8, color=MID_GREY, align=PP_ALIGN.LEFT)

def slide_bg(slide, color=LIGHT_GREY):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / COVER
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SAP_DARK)
# diagonal accent
add_rect(s, 0, 4.8, 13.33, 0.06, fill=SAP_BLUE)
add_rect(s, 0, 4.86, 13.33, 2.64, fill=RGBColor(0x00, 0x1F, 0x44))

# tank icon placeholder (circle)
add_rect(s, 10.8, 0.4, 2.0, 3.8, fill=RGBColor(0x00, 0x1F, 0x44), line=SAP_BLUE, line_w=1.5)
add_text(s, "🛢", 11.1, 1.4, 1.4, 1.4, size=54, align=PP_ALIGN.CENTER, color=WHITE)

add_text(s, "Hydrocarbon Tank Stock", 0.6, 1.0, 9.8, 0.9,
         size=38, bold=True, color=WHITE)
add_text(s, "Reconciliation", 0.6, 1.85, 9.8, 0.8,
         size=38, bold=True, color=ACCENT_GOLD)
add_text(s, "End-to-End Automated Pipeline  |  SAP BTP + IS-OIL HPM on S/4HANA Private Cloud",
         0.6, 2.75, 9.8, 0.5, size=14, color=MID_GREY)

add_rect(s, 0.6, 3.4, 5.5, 0.05, fill=SAP_BLUE)

add_text(s, "From dip reading to Material Document in under 2 minutes",
         0.6, 3.6, 9.8, 0.5, size=13, italic=True, color=RGBColor(0x88, 0xCC, 0xFF))

# bottom strip badges
badges = [("SAP BTP", SAP_BLUE), ("IS-OIL HPM", SAP_TEAL),
          ("CAP + React", SAP_GREEN), ("n8n Workflow", SAP_AMBER),
          ("Cloud Connector", SAP_RED)]
bx = 0.6
for label, col in badges:
    add_rect(s, bx, 5.1, 1.8, 0.45, fill=col)
    add_text(s, label, bx+0.05, 5.15, 1.7, 0.35,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bx += 2.0

add_text(s, "Terminal Operations / Supply Chain", 0.6, 5.9, 8, 0.4,
         size=11, color=MID_GREY)
add_text(s, "2026", 0.6, 6.3, 3, 0.4, size=11, color=MID_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA / TABLE OF CONTENTS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)
header_bar(s, "Agenda", "What we will cover")
footer_bar(s)

items = [
    ("01", "Business Problem",          "Manual, fragmented, error-prone reconciliation"),
    ("02", "Solution Overview",         "Automated pipeline: gauge reading → SAP posting"),
    ("03", "Architecture",              "IS-OIL OGS/650 | SAP BTP | React Dashboard"),
    ("04", "End-to-End Flowchart",      "6 milestones from data ingestion to report"),
    ("05", "Use Cases & User Roles",    "5 personas: Operator → Supervisor → Finance"),
    ("06", "Intelligent Features",      "AI classifications, alerts, state machine, audit"),
    ("07", "Key Achievements",          "2 hrs → 2 min  •  First IS-OIL↔BTP integration"),
    ("08", "Technology Stack",          "Components, APIs, and integration points"),
]

cols = [0.4, 6.9]
for i, (num, title, desc) in enumerate(items):
    col = i % 2
    row = i // 2
    x = cols[col]
    y = 1.55 + row * 1.28
    add_rect(s, x, y, 6.0, 1.1, fill=LIGHT_GREY, line=MID_GREY, line_w=0.5)
    add_rect(s, x, y, 0.55, 1.1, fill=SAP_BLUE)
    add_text(s, num, x+0.02, y+0.28, 0.52, 0.5,
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, x+0.65, y+0.08, 5.2, 0.4,
             size=13, bold=True, color=SAP_DARK)
    add_text(s, desc, x+0.65, y+0.52, 5.2, 0.45,
             size=10, color=RGBColor(0x55, 0x55, 0x55), italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BUSINESS PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)
header_bar(s, "The Business Problem", "Manual reconciliation — fragmented, slow, error-prone")
footer_bar(s)

# Left panel — problem statement
add_rect(s, 0.35, 1.45, 5.8, 5.55, fill=LIGHT_GREY, line=MID_GREY, line_w=0.5)
add_rect(s, 0.35, 1.45, 5.8, 0.45, fill=SAP_RED)
add_text(s, "The Manual Process Today", 0.55, 1.5, 5.4, 0.38,
         size=12, bold=True, color=WHITE)

steps = [
    "1  Operator physically dips tank or reads ATG gauge",
    "2  Writes reading on paper / spreadsheet",
    "3  Logs into SAP to look up book stock manually",
    "4  Calculates delta using calculator or Excel",
    "5  Applies VCF correction from printed lookup tables",
    "6  Classifies variance from memory (OK/FLAG/URGENT)",
    "7  Emails supervisor for URGENT approval",
    "8  Supervisor replies by email — or doesn't (delays!)",
    "9  Manually posts goods movement in SAP MIGO",
    "10  Updates shared spreadsheet, emails report to Finance",
]
for i, step in enumerate(steps):
    y = 2.05 + i * 0.45
    add_rect(s, 0.45, y, 5.6, 0.38,
             fill=WHITE if i % 2 == 0 else LIGHT_GREY)
    add_text(s, step, 0.6, y+0.04, 5.3, 0.32, size=9.5, color=TEXT_DARK)

# Right panel — consequences
add_rect(s, 6.6, 1.45, 6.38, 2.55, fill=RGBColor(0xFF, 0xF0, 0xF0), line=SAP_RED, line_w=0.8)
add_rect(s, 6.6, 1.45, 6.38, 0.45, fill=SAP_RED)
add_text(s, "Consequences", 6.8, 1.5, 6.0, 0.38, size=12, bold=True, color=WHITE)

consequences = [
    ("⏱", "2–4 hours per terminal per day", SAP_RED),
    ("❌", "Error-prone — wrong VCF table, wrong material", SAP_RED),
    ("🔓", "No audit trail — email chains are not tamper-evident", SAP_AMBER),
    ("🐢", "URGENT variances sit unresolved for hours", SAP_AMBER),
    ("📊", "Data split across ATG, SAP, spreadsheets, email", SAP_DARK),
]
for i, (icon, text, col) in enumerate(consequences):
    y = 2.05 + i * 0.43
    add_text(s, icon, 6.7, y, 0.45, 0.38, size=13, align=PP_ALIGN.CENTER)
    add_rect(s, 7.15, y+0.06, 5.6, 0.28, fill=col)
    add_text(s, text, 7.2, y+0.07, 5.5, 0.26, size=9.5, bold=False, color=WHITE)

# Right panel — hidden complexity
add_rect(s, 6.6, 4.2, 6.38, 2.8, fill=RGBColor(0xE8, 0xF0, 0xF8), line=SAP_BLUE, line_w=0.8)
add_rect(s, 6.6, 4.2, 6.38, 0.45, fill=SAP_BLUE)
add_text(s, "Hidden Complexity — IS-OIL", 6.8, 4.25, 6.0, 0.38, size=12, bold=True, color=WHITE)
hc = [
    "Standard SAP has no concept of a 'tank'",
    "IS-OIL HPM extends S/4HANA with tank master data,",
    "  dip history, strapping tables & VCF conversion",
    "None of this is exposed via standard public OData APIs",
    "Requires deep ABAP knowledge of IS-OIL internals",
]
for i, line in enumerate(hc):
    add_text(s, ("• " if not line.startswith(" ") else "") + line.strip(),
             6.75, 4.82 + i*0.42, 6.1, 0.38, size=9.5, color=SAP_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION OVERVIEW  (Before vs After)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)
header_bar(s, "Solution Overview", "What changed — before vs. after")
footer_bar(s)

before_after = [
    ("Operator reads dip manually, writes on paper",
     "IS-OIL OIB_TANKDIP read automatically via custom OData service"),
    ("Book stock looked up manually in SAP",
     "RELSTOCK fetched directly from IS-OIL dip record"),
    ("VCF correction done with Excel/lookup tables",
     "VCF Calculator applies automatically; ASTM fallback if API unavailable"),
    ("Delta calculated in spreadsheet",
     "Variance Engine computes delta and classifies in milliseconds"),
    ("Supervisor approval via email",
     "Approval Queue in CAP dashboard with full audit record"),
    ("Goods movement posted manually in MIGO",
     "Material Document created automatically via API_MATERIAL_DOCUMENT_SRV"),
    ("Report emailed manually",
     "PDF generated and distributed to Email + MS Teams automatically"),
    ("No audit trail",
     "Every milestone, decision & posting recorded with timestamp & actor"),
]

# Column headers
add_rect(s, 0.35, 1.42, 5.85, 0.45, fill=SAP_RED)
add_text(s, "BEFORE  (Manual)", 0.5, 1.47, 5.6, 0.38, size=12, bold=True, color=WHITE)
add_rect(s, 6.55, 1.42, 6.38, 0.45, fill=SAP_GREEN)
add_text(s, "AFTER  (Automated)", 6.7, 1.47, 6.1, 0.38, size=12, bold=True, color=WHITE)
add_rect(s, 6.3, 1.42, 0.25, 0.45, fill=ACCENT_GOLD)

for i, (bef, aft) in enumerate(before_after):
    y = 1.95 + i * 0.62
    bg = WHITE if i % 2 == 0 else LIGHT_GREY
    add_rect(s, 0.35, y, 5.85, 0.55, fill=bg)
    add_rect(s, 6.55, y, 6.38, 0.55, fill=bg)
    add_rect(s, 6.3, y, 0.25, 0.55, fill=ACCENT_GOLD)
    add_text(s, "✗  " + bef, 0.5, y+0.06, 5.6, 0.44, size=9.5, color=SAP_RED)
    add_text(s, "✓  " + aft, 6.68, y+0.06, 6.1, 0.44, size=9.5, color=SAP_GREEN)

# KPI strip at bottom
add_rect(s, 0.35, 6.82, 12.63, 0.0)  # spacer
kpis = [
    ("2–4 hrs → <2 min", "Daily reconciliation time"),
    ("Zero manual steps", "For within-tolerance tanks"),
    ("100% audit coverage", "Every action timestamped"),
    ("Real-time alerts", "OK / FLAG / URGENT"),
]
for i, (val, label) in enumerate(kpis):
    x = 0.35 + i * 3.2
    add_rect(s, x, 6.7, 3.0, 0.6, fill=SAP_DARK)
    add_text(s, val, x+0.1, 6.72, 2.8, 0.3, size=11, bold=True,
             color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    add_text(s, label, x+0.1, 6.99, 2.8, 0.25, size=8.5, color=WHITE,
             align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)
header_bar(s, "Solution Architecture", "Three layers — IS-OIL OGS/650 | SAP BTP | React Dashboard")
footer_bar(s)

# Layer boxes
layers = [
    (0.3,  "LAYER 1 — Field / IS-OIL OGS/650 (Source of Truth)", SAP_RED,
     ["OIB_TANKDIP table — live tank dip readings",
      "ZTANK_DIP_SRV_SRV — custom OData (TankDipSet, ReasonCodeSet, NominationSet)",
      "ZTANK_PLANT_SRV_SRV — plant/terminal list",
      "ZTANK_POST_SRV_SRV — goods movement posting via IS-OIL BAPIs",
      "Z_TANK_RECON_TRIGGER_RUN — ABAP FM for OGS → BTP M2M trigger"]),
    (2.85, "LAYER 2 — SAP BTP (Automation Engine)", SAP_BLUE,
     ["n8n Workflow — orchestrates all 6 milestones end-to-end",
      "CAP Node.js Application — persistent backbone: runs, variances, approvals, audit log",
      "XSUAA — role-based access control (4 roles: User / Approver / Admin / OGS)",
      "BTP Alert Notification Service — push alerts to Finance & Supervisor",
      "Cloud Connector (APAC_DEV10) — secure tunnel to on-premise OGS/650"]),
    (5.4,  "LAYER 3 — React Dashboard (Single Cockpit)", SAP_GREEN,
     ["Dashboard — live run status & per-tank results (all roles)",
      "Approval Queue — URGENT variance review (Supervisor only)",
      "Tank Detail — VCF-corrected volumes, delta, audit drill-down",
      "Audit Trail — complete M1→M6 history with timestamps & actors",
      "Trend Chart — 30-day variance history per tank",
      "AI Chat — natural-language query interface"]),
]

for (y, title, col, points) in layers:
    add_rect(s, 0.25, y+1.38, 12.83, 2.3, fill=LIGHT_GREY, line=col, line_w=1.2)
    add_rect(s, 0.25, y+1.38, 12.83, 0.42, fill=col)
    add_text(s, title, 0.45, y+1.42, 12.4, 0.35, size=11, bold=True, color=WHITE)
    cols_n = 2
    per_col = (len(points) + 1) // 2
    for j, pt in enumerate(points):
        cx = j // per_col
        cy = j % per_col
        add_text(s, "▸  " + pt,
                 0.4 + cx * 6.5, y+1.92 + cy * 0.38,
                 6.2, 0.35, size=9.5, color=TEXT_DARK)

# Connector arrows
for y_arrow in [3.68, 5.68]:
    add_rect(s, 5.9, y_arrow, 1.55, 0.35, fill=ACCENT_GOLD)
    add_text(s, "Cloud Connector  ▼" if y_arrow < 4 else "OData / REST  ▼",
             5.92, y_arrow+0.04, 1.5, 0.28, size=8.5, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

# Right-side integration callout
add_rect(s, 9.8, 2.2, 3.2, 4.75, fill=RGBColor(0xE8, 0xF0, 0xF8), line=SAP_TEAL, line_w=0.8)
add_rect(s, 9.8, 2.2, 3.2, 0.38, fill=SAP_TEAL)
add_text(s, "S/4HANA APIs", 9.9, 2.23, 3.0, 0.32, size=10, bold=True, color=WHITE)
apis = [
    "API_MATERIAL_STOCK_SRV",
    "(HPM Book Stock)",
    "",
    "API_PHYSICAL_INVENTORY_DOC_SRV",
    "(Fiori Dip Entries)",
    "",
    "MEASUREMENTDOCUMENT_0001",
    "(Tank Master + Strapping)",
    "",
    "API_MATERIAL_DOCUMENT_SRV",
    "(Goods Movement POST)",
    "",
    "Hydrocarbon Qty Conv. REST",
    "(VCF Calculation)",
]
for i, api in enumerate(apis):
    col = SAP_BLUE if not api.startswith("(") and api != "" else RGBColor(0x44, 0x66, 0x88)
    bold = not api.startswith("(") and api != ""
    add_text(s, api, 9.9, 2.68 + i * 0.29, 3.05, 0.28,
             size=8 if api.startswith("(") else 8.5, bold=bold, color=col, italic=api.startswith("("))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — END-TO-END FLOWCHART
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)
header_bar(s, "End-to-End Flow — 6 Milestones", "From tank gauge reading to Material Document and PDF report")
footer_bar(s)

milestones = [
    ("M1", "Data\nIngestion", SAP_TEAL,
     "ATG + Fiori dip readings\nRead OIB_TANKDIP via\nZTANK_DIP_SRV_SRV\nValidate completeness"),
    ("M2", "VCF\nCorrection", SAP_BLUE,
     "Gross → Net volume\nVia Hydrocarbon Qty\nConversion REST API\nASTM fallback if needed"),
    ("M3", "Variance\nCalculation", SAP_AMBER,
     "Delta = Physical − Book\nClassify:\n🟢 OK ≤0.10%\n🟡 FLAG ≤0.25%\n🔴 URGENT >0.25%"),
    ("M4", "Approval\nDecision", SAP_RED,
     "URGENT held in CAP\nSupervisor reviews\nApprove / Reject\nReason code mandatory"),
    ("M5", "Goods Movement\nPosting", SAP_GREEN,
     "Create → Count → PostDiff\nvia IS-OIL BAPIs\nMaterial Doc hyperlink\nAI rec. on failure"),
    ("M6", "Report &\nAlert", SAP_DARK,
     "PDF per-tank variance\nEmail + MS Teams\nBTP Alert Notification\nAMBER auto-post @8hrs"),
]

# Draw flow row
box_w, box_h = 1.82, 1.25
box_y = 1.55
gap = 0.2
start_x = 0.35

for i, (m, title, col, desc) in enumerate(milestones):
    bx = start_x + i * (box_w + gap)
    # main box
    add_rect(s, bx, box_y, box_w, box_h, fill=col)
    add_text(s, m, bx+0.04, box_y+0.06, 0.45, 0.35,
             size=13, bold=True, color=WHITE)
    add_text(s, title, bx+0.02, box_y+0.38, box_w-0.08, 0.8,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # arrow between boxes
    if i < 5:
        ax = bx + box_w
        add_rect(s, ax, box_y + box_h/2 - 0.07, gap, 0.14, fill=ACCENT_GOLD)

# Description cards below
for i, (m, title, col, desc) in enumerate(milestones):
    bx = start_x + i * (box_w + gap)
    cy = box_y + box_h + 0.12
    add_rect(s, bx, cy, box_w, 2.75, fill=LIGHT_GREY, line=col, line_w=1.0)
    # connector line
    add_rect(s, bx + box_w/2 - 0.04, cy - 0.12, 0.08, 0.14, fill=col)
    for j, line in enumerate(desc.split("\n")):
        add_text(s, line, bx+0.08, cy+0.12+j*0.47, box_w-0.16, 0.44,
                 size=9, color=TEXT_DARK if j > 0 else col,
                 bold=(j == 0))

# Decision diamond for M3→M4 branch
add_rect(s, 6.62, 4.55, 2.1, 0.6, fill=RGBColor(0xFF, 0xF0, 0xD0), line=SAP_AMBER, line_w=0.8)
add_text(s, "URGENT?", 6.67, 4.6, 2.0, 0.48, size=9.5, bold=True,
         color=SAP_AMBER, align=PP_ALIGN.CENTER)
add_text(s, "YES →", 6.67, 5.1, 0.7, 0.25, size=8, color=SAP_RED, bold=True)
add_text(s, "NO → auto-post at M5", 7.3, 5.1, 1.8, 0.25, size=8, color=SAP_GREEN)

# Bottom legend
legend = [("🟢 OK", "Within tolerance — auto-post", SAP_GREEN),
          ("🟡 FLAG", "Exceeds 2nd threshold — auto-post + notify", SAP_AMBER),
          ("🔴 URGENT", "Exceeds primary — HOLD for supervisor", SAP_RED)]
lx = 0.35
for lbl, desc, col in legend:
    add_rect(s, lx, 7.02, 4.15, 0.32, fill=col)
    add_text(s, f"{lbl}  {desc}", lx+0.1, 7.04, 4.0, 0.26,
             size=8.5, bold=False, color=WHITE)
    lx += 4.3


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — USE CASES & USER ROLES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)
header_bar(s, "Use Cases & User Roles", "Five personas — one cockpit")
footer_bar(s)

roles = [
    ("Terminal\nOperator", "Field / Gauger", SAP_TEAL,
     ["Enter manual dip readings via Fiori app",
      "View ingestion status in dashboard",
      "See today's tank readings & run status"]),
    ("Stock\nController", "Inventory Analyst", SAP_BLUE,
     ["Monitor daily run: Ingestion→VCF→Variance→Posting",
      "Drill into per-tank VCF-corrected detail",
      "Re-trigger data collection on partial failure",
      "Download variance PDF for any completed run"]),
    ("Terminal\nSupervisor", "Carlos — Operations Lead", SAP_RED,
     ["Receive URGENT alert immediately",
      "Review URGENT-flagged tank detail in approval queue",
      "Approve / Reject with reason code (mandatory on reject)",
      "Receive consolidated PDF report in MS Teams"]),
    ("Finance\nAccountant", "Fatima — Finance Officer", SAP_AMBER,
     ["Receive OK/FLAG/URGENT alerts via BTP ANS",
      "View posting status and Material Document IDs",
      "Access audit trail for daily inventory valuation",
      "No dashboard operation needed — trust the alert"]),
    ("Compliance\nOfficer", "Read-Only", SAP_DARK,
     ["Browse full historical audit trail",
      "Filter by tank, date, classification",
      "Verify audit records for regulatory submissions"]),
]

rw, rh = 2.35, 5.2
rx = 0.28
for i, (role, sub, col, tasks) in enumerate(roles):
    x = rx + i * (rw + 0.12)
    # header
    add_rect(s, x, 1.45, rw, 1.0, fill=col)
    add_text(s, role, x+0.08, 1.5, rw-0.16, 0.62,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, sub, x+0.08, 2.07, rw-0.16, 0.32,
             size=8.5, italic=True, color=RGBColor(0xDD,0xDD,0xDD),
             align=PP_ALIGN.CENTER)
    # task list
    add_rect(s, x, 2.45, rw, 3.85, fill=LIGHT_GREY, line=col, line_w=0.6)
    for j, task in enumerate(tasks):
        ty = 2.58 + j * 0.82
        add_rect(s, x+0.1, ty, rw-0.2, 0.7, fill=WHITE)
        add_rect(s, x+0.1, ty, 0.06, 0.7, fill=col)
        add_text(s, task, x+0.22, ty+0.08, rw-0.35, 0.56,
                 size=9, color=TEXT_DARK)

# RBAC note at bottom
add_rect(s, 0.28, 6.4, 12.77, 0.52, fill=SAP_DARK)
add_text(s,
    "Role-Based Access Control enforced via XSUAA scopes:  "
    "ReconciliationUser (read)  |  ReconciliationApprover (approve/reject)  |  "
    "ReconciliationAdmin (config + trigger)  |  OGSIntegration (M2M)",
    0.45, 6.44, 12.4, 0.42, size=9, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — INTELLIGENT FEATURES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)
header_bar(s, "Intelligent Features", "AI-powered automation built into every milestone")
footer_bar(s)

features = [
    ("AI Variance\nClassification", SAP_BLUE,
     "Variance Engine auto-classifies each tank as\n"
     "GREEN / AMBER / RED against configurable thresholds.\n"
     "No manual judgment — exceptions surface automatically."),
    ("AI Recommendations\non Failure", SAP_RED,
     "When goods movement posting fails (M5),\n"
     "AI generates corrective action recommendations\n"
     "surfaced in audit trail and dashboard."),
    ("Intelligent Approval\nState Machine", SAP_TEAL,
     "CAP holds URGENT tanks; auto-approves GREEN/AMBER.\n"
     "Human involvement limited to genuine exceptions.\n"
     "Approval record mandatory before any posting fires."),
    ("Time-Aware\nAuto-Post Rule", SAP_AMBER,
     "AMBER-classified tanks auto-post after 8 hours\n"
     "if no supervisor action — prevents reconciliation\n"
     "backlog without sacrificing oversight."),
    ("Dynamic Reason\nCodes (M4)", SAP_GREEN,
     "Reason codes fetched live from SAP table T157D —\n"
     "not hardcoded. Adapts to master data configuration\n"
     "in the underlying S/4HANA system."),
    ("Intelligent Alert\nRouting", SAP_DARK,
     "Alert Manager routes OK/FLAG/URGENT notifications\n"
     "to the right roles via the right channels\n"
     "(Teams webhook, BTP ANS, in-app bell) by severity."),
    ("Dual-Source Data\nValidation", SAP_BLUE,
     "Data Collector merges ATG + Fiori manual dip entries\n"
     "and cross-validates for completeness.\n"
     "Rejects incomplete runs before any calculation."),
    ("AI Chat Interface\n(Natural Language)", SAP_TEAL,
     "Operators can query reconciliation data in plain English\n"
     "via AI Core proxy — 'What is Tank 17T1 variance today?'\n"
     "Surfaced as an AiChat page in the React dashboard."),
    ("VCF Fallback\nIntelligence", SAP_AMBER,
     "If SAP Hydrocarbon Qty Conversion API is unavailable,\n"
     "ASTM D1250 fallback table activates automatically.\n"
     "Fallback flagged in audit trail and dashboard UI."),
]

cols_n = 3
fw, fh = 4.0, 1.88
fx0 = 0.32
fy0 = 1.52
gap_x, gap_y = 0.15, 0.16

for i, (title, col, desc) in enumerate(features):
    c = i % cols_n
    r = i // cols_n
    fx = fx0 + c * (fw + gap_x)
    fy = fy0 + r * (fh + gap_y)
    add_rect(s, fx, fy, fw, fh, fill=LIGHT_GREY, line=col, line_w=0.8)
    add_rect(s, fx, fy, fw, 0.42, fill=col)
    add_text(s, title, fx+0.08, fy+0.04, fw-0.16, 0.36,
             size=10.5, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(s, desc, fx+0.1, fy+0.5, fw-0.2, 1.3,
             size=9, color=TEXT_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)
header_bar(s, "Technology Stack", "Components, APIs, and integration points")
footer_bar(s)

# Left: component table
add_rect(s, 0.3, 1.45, 7.8, 0.42, fill=SAP_DARK)
add_text(s, "Component", 0.5, 1.5, 2.6, 0.34, size=10, bold=True, color=WHITE)
add_text(s, "Technology", 3.15, 1.5, 2.5, 0.34, size=10, bold=True, color=WHITE)
add_text(s, "Purpose", 5.7, 1.5, 2.3, 0.34, size=10, bold=True, color=WHITE)

components = [
    ("Scheduler",          "n8n Cron trigger",         "Daily run at configured time"),
    ("Data Collector",     "n8n workflow node",         "ATG + Fiori ingestion & validation"),
    ("VCF Calculator",     "n8n workflow node",         "Gross → Net volume conversion"),
    ("Variance Engine",    "n8n workflow node",         "Delta compute & OK/FLAG/URGENT"),
    ("Alert Manager",      "n8n workflow node",         "Audit log + BTP ANS alerts"),
    ("Report Generator",   "n8n workflow node",         "PDF + Email + MS Teams distrib."),
    ("CAP Backend",        "SAP CAP Node.js / BTP",     "Persistent store + approval state"),
    ("React Dashboard",    "React 19 + SAP UI5 WC",     "Single cockpit for all roles"),
    ("Auth",               "XSUAA (BTP)",               "Role-based access control"),
    ("OData Services",     "ABAP Gateway (OGS/650)",    "Custom IS-OIL exposure layer"),
    ("Connectivity",       "SAP Cloud Connector",       "Secure on-premise tunnel"),
]

for i, (comp, tech, purpose) in enumerate(components):
    y = 1.9 + i * 0.44
    bg = WHITE if i % 2 == 0 else LIGHT_GREY
    add_rect(s, 0.3, y, 7.8, 0.42, fill=bg)
    add_text(s, comp,    0.45, y+0.06, 2.6, 0.32, size=9.5, bold=True,  color=SAP_DARK)
    add_text(s, tech,    3.15, y+0.06, 2.4, 0.32, size=9,   color=SAP_BLUE)
    add_text(s, purpose, 5.7,  y+0.06, 2.3, 0.32, size=9,   color=TEXT_DARK)

# Right: API table
add_rect(s, 8.4, 1.45, 4.6, 0.42, fill=SAP_TEAL)
add_text(s, "S/4HANA / OGS APIs Used", 8.55, 1.5, 4.3, 0.34,
         size=10, bold=True, color=WHITE)

apis_list = [
    ("ZTANK_DIP_SRV_SRV",                "OGS — Tank dip readings",     "Read"),
    ("ZTANK_PLANT_SRV_SRV",              "OGS — Plant / terminal list", "Read"),
    ("ZTANK_POST_SRV_SRV",               "OGS — Goods movement post",   "Write"),
    ("API_MATERIAL_STOCK_SRV",           "S/4 — HPM book stock",        "Read"),
    ("API_PHYSICAL_INVENTORY_DOC_SRV",   "S/4 — Fiori dip entries",     "Read"),
    ("MEASUREMENTDOCUMENT_0001",         "S/4 — Tank master/strapping", "Read"),
    ("Hydrocarbon Qty Conv. REST API",   "S/4 — VCF calculation",       "Read"),
    ("API_MATERIAL_DOCUMENT_SRV",        "S/4 — Goods movement POST",   "Write ⚠"),
    ("BTP Alert Notification Service",  "BTP — Push alerts",           "Write"),
    ("MS Teams webhook",                 "Teams — PDF report distrib.", "Write"),
]

for i, (api, desc, rw_flag) in enumerate(apis_list):
    y = 1.9 + i * 0.44
    bg = WHITE if i % 2 == 0 else LIGHT_GREY
    add_rect(s, 8.4, y, 4.6, 0.42, fill=bg)
    flag_col = SAP_RED if "Write" in rw_flag else SAP_GREEN
    add_rect(s, 8.4, y, 0.6, 0.42, fill=flag_col)
    add_text(s, rw_flag.replace(" ⚠",""), 8.41, y+0.07, 0.58, 0.28,
             size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, api,  9.05, y+0.06, 2.4, 0.18, size=7.8, bold=True, color=SAP_DARK)
    add_text(s, desc, 9.05, y+0.23, 2.4, 0.18, size=7.5, color=RGBColor(0x55,0x55,0x55), italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — KEY ACHIEVEMENTS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)
header_bar(s, "Key Achievements", "Outcomes delivered — technical and business")
footer_bar(s)

achievements = [
    ("First IS-OIL → BTP\nIntegration", SAP_BLUE,
     "Connected SAP BTP to IS-OIL HPM tank data on a\n"
     "Private Cloud system via Cloud Connector.\n"
     "Real OIB_TANKDIP data verified against SE16N records."),
    ("Custom OData Services\nBuilt from Scratch", SAP_TEAL,
     "ZTANK_DIP_SRV_SRV (3 entity sets),\n"
     "ZTANK_PLANT_SRV_SRV, ZTANK_POST_SRV_SRV —\n"
     "IS-OIL tank data exposed as OData for the first time."),
    ("OGS → BTP M2M\nIntegration", SAP_GREEN,
     "ABAP FM Z_TANK_RECON_TRIGGER_RUN fetches\n"
     "XSUAA token and calls BTP CAP endpoint directly —\n"
     "IS-OIL can trigger reconciliation runs from ABAP."),
    ("Full Approval Governance\non BTP", SAP_AMBER,
     "URGENT variances held in CAP state machine.\n"
     "No Material Document posted without\n"
     "supervisor sign-off recorded in immutable audit trail."),
    ("End-to-End Data\nAccuracy Verified", SAP_RED,
     "Dashboard values match OIB_TANKDIP source exactly.\n"
     "Book stock, physical quantity, and delta all verified\n"
     "against SE16N records in OGS/650."),
    ("Role-Based Access\nControl", SAP_DARK,
     "4 XSUAA scopes enforced at API and UI levels.\n"
     "Approval Queue visible to Supervisors only.\n"
     "Config screen accessible to Admins only."),
]

aw, ah = 3.95, 2.15
ax0 = 0.3
ay0 = 1.52
for i, (title, col, desc) in enumerate(achievements):
    c = i % 3
    r = i // 3
    ax = ax0 + c * (aw + 0.15)
    ay = ay0 + r * (ah + 0.18)
    add_rect(s, ax, ay, aw, ah, fill=LIGHT_GREY, line=col, line_w=1.0)
    add_rect(s, ax, ay, aw, 0.5, fill=col)
    add_text(s, "✅", ax+0.08, ay+0.07, 0.4, 0.38, size=14, color=WHITE)
    add_text(s, title, ax+0.48, ay+0.08, aw-0.6, 0.4,
             size=10.5, bold=True, color=WHITE)
    add_text(s, desc, ax+0.12, ay+0.62, aw-0.24, 1.42,
             size=9, color=TEXT_DARK)

# Hero KPI banner
add_rect(s, 0.3, 6.65, 12.73, 0.62, fill=SAP_DARK)
add_text(s,
    "🏆   2–4 hours of daily manual work  →  under 2 minutes automated  "
    "  |  First live IS-OIL ↔ BTP integration in this landscape  "
    "  |  100% audit coverage, zero manual steps for within-tolerance tanks",
    0.5, 6.7, 12.3, 0.5, size=10.5, bold=True, color=ACCENT_GOLD,
    align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RISKS & GOVERNANCE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)
header_bar(s, "Risks, Guardrails & Governance", "What can go wrong — and how we handle it")
footer_bar(s)

# Left: risks
add_rect(s, 0.3, 1.45, 6.1, 0.42, fill=SAP_RED)
add_text(s, "Key Risks", 0.5, 1.5, 5.8, 0.34, size=11, bold=True, color=WHITE)

risks = [
    ("HIGH", "Tank Strapping Table OData Coverage",
     "Standard OData may not fully expose strapping calibration tables.\n"
     "RFC/BAPI extension may be required. Validate in Sprint 1."),
    ("MED", "ATG Protocol Variability",
     "ATG vendors differ in payload format. Data Collector normalisation\n"
     "layer must be designed for pluggability across sites."),
    ("MED", "VCF REST API Stability",
     "Hydrocarbon Qty Conversion REST API has limited documentation.\n"
     "ASTM fallback must be operational from day one."),
    ("MED", "S/4HANA Auth Provisioning",
     "Material Document posting OData calls require specific auth profiles.\n"
     "Must be provisioned and tested before UAT."),
]
for i, (level, title, desc) in enumerate(risks):
    y = 1.94 + i * 1.18
    col = SAP_RED if level == "HIGH" else SAP_AMBER
    add_rect(s, 0.3, y, 6.1, 1.1, fill=LIGHT_GREY, line=col, line_w=0.8)
    add_rect(s, 0.3, y, 0.62, 1.1, fill=col)
    add_text(s, level, 0.31, y+0.35, 0.6, 0.35, size=9, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, 1.0, y+0.1, 5.25, 0.32, size=10, bold=True, color=col)
    add_text(s, desc, 1.0, y+0.44, 5.25, 0.58, size=9, color=TEXT_DARK)

# Right: guardrails
add_rect(s, 6.75, 1.45, 6.22, 0.42, fill=SAP_GREEN)
add_text(s, "System Guardrails & Fail-safes", 6.9, 1.5, 5.9, 0.34, size=11, bold=True, color=WHITE)

guardrails = [
    ("🔒", "No URGENT posting without supervisor approval record in CAP"),
    ("🛑", "Incomplete ingestion halts run before VCF — URGENT alert raised"),
    ("⚠", "Posting failure → logged + URGENT alert; no auto-retry"),
    ("🔄", "VCF API unavailable → ASTM fallback auto-activates; flagged in audit"),
    ("📋", "Every action timestamped with actor ID — tamper-evident audit log"),
    ("🔑", "XSUAA scopes enforced at API and UI level for every role"),
    ("⏱", "AMBER auto-posts after 8 hours — prevents reconciliation backlog"),
    ("🔗", "Material Doc number written back to CAP for full traceability"),
]
for i, (icon, text) in enumerate(guardrails):
    y = 1.94 + i * 0.6
    add_rect(s, 6.75, y, 6.22, 0.54, fill=WHITE if i % 2 == 0 else LIGHT_GREY)
    add_text(s, icon, 6.82, y+0.1, 0.42, 0.35, size=14, align=PP_ALIGN.CENTER)
    add_text(s, text, 7.3, y+0.1, 5.55, 0.35, size=10, color=TEXT_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BUSINESS IMPACT  (crisp version)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, WHITE)
header_bar(s, "Business Impact", "What this solution changes — in numbers and in practice")
footer_bar(s)

# ── Hero metric banner (full width, 5 numbers) ────────────────────────────
metrics = [
    ("~95%",      "Reduction in\nreconciliation time"),
    ("< 2 min",   "Gauge reading\nto Material Document"),
    ("100%",      "Auto-posted\n(within-tolerance tanks)"),
    ("0",         "Manual steps\nfor routine runs"),
    ("6",         "Milestones fully\naudit-logged"),
]
for i, (num, lbl) in enumerate(metrics):
    x = 0.28 + i * 2.56
    add_rect(s, x, 1.44, 2.4, 1.05, fill=SAP_DARK)
    add_text(s, num, x+0.1, 1.48, 2.2, 0.52,
             size=26, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x+0.1, 1.97, 2.2, 0.46,
             size=8.5, color=MID_GREY, align=PP_ALIGN.CENTER)

# ── Three crisp impact columns ────────────────────────────────────────────
cols_data = [
    ("Operational", SAP_BLUE, [
        ("⏱  2–4 hrs → < 2 min",
         "Daily reconciliation time per terminal",
         "n8n auto-runs ATG ingestion, VCF correction, and posting — no one touches it"),
        ("🖥  One cockpit, five roles",
         "No more SAP screens + spreadsheets + email",
         "Every role sees live run status, tank detail, and approvals in one React dashboard"),
        ("🤖  Zero manual steps",
         "ATG → VCF → post runs fully unattended",
         "Routine tanks post automatically; humans only act on URGENT exceptions"),
        ("📈  Scales to new terminals",
         "Config only — no extra process or headcount",
         "Add terminal in CAP configuration — the same pipeline runs for every site"),
    ]),
    ("Financial", SAP_GREEN, [
        ("💰  Same-day inventory close",
         "Material Docs posted minutes after run — not hours",
         "Finance gets accurate stock values the same morning, not after chasing approvals"),
        ("🔍  No missed shrinkage",
         "Every tank variance caught and classified automatically",
         "Variance Engine checks all tanks every run — nothing slips through a manual gap"),
        ("✅  Correct move type always",
         "551/552 determined by engine — no manual posting error",
         "Shrinkage vs gain selected by delta sign — eliminates wrong-type reversals"),
        ("📄  Audit-ready records",
         "Every Doc linked to VCF factor, approval actor, timestamp",
         "Auditors get the full posting chain from one CAP query — no email reconstruction"),
    ]),
    ("Compliance & Risk", SAP_TEAL, [
        ("🔒  Structurally enforced",
         "URGENT tanks cannot post without supervisor sign-off",
         "CAP state machine blocks the API call — approval is not a reminder, it is a gate"),
        ("📋  Immutable audit trail",
         "Full M1–M6 chain of custody — every action timestamped",
         "Every reading, VCF factor, decision, and posting stored with actor ID and timestamp"),
        ("🔎  Instant regulatory lookup",
         "Filter by tank / date / class — reconstruct any run in seconds",
         "Compliance officers query the audit trail directly — no spreadsheet archaeology"),
        ("📐  VCF fully traceable",
         "Strapping ID + temp factor logged — meets custody transfer reqs",
         "Each corrected volume cites its exact calibration source — defensible in any audit"),
    ]),
]

cw, ch = 4.1, 4.75
cx0, cy0 = 0.28, 2.62
for ci, (title, col, rows) in enumerate(cols_data):
    cx = cx0 + ci * (cw + 0.165)
    add_rect(s, cx, cy0, cw, ch, fill=LIGHT_GREY, line=col, line_w=1.1)
    add_rect(s, cx, cy0, cw, 0.42, fill=col)
    add_text(s, title, cx+0.12, cy0+0.06, cw-0.24, 0.33, size=12, bold=True, color=WHITE)
    for j, (headline, subline, how) in enumerate(rows):
        ry = cy0 + 0.5 + j * 1.06
        add_rect(s, cx+0.1, ry, cw-0.2, 0.98, fill=WHITE)
        add_rect(s, cx+0.1, ry, 0.05, 0.98, fill=col)
        add_text(s, headline, cx+0.22, ry+0.04, cw-0.36, 0.28,
                 size=10, bold=True, color=SAP_DARK)
        add_text(s, subline,  cx+0.22, ry+0.32, cw-0.36, 0.26,
                 size=8.5, color=col, italic=False, bold=False)
        add_text(s, how,      cx+0.22, ry+0.57, cw-0.36, 0.36,
                 size=8, color=RGBColor(0x44, 0x44, 0x55), italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CLOSING / SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=SAP_DARK)
add_rect(s, 0, 5.1, 13.33, 0.06, fill=SAP_BLUE)
add_rect(s, 0, 5.16, 13.33, 2.34, fill=RGBColor(0x00, 0x1F, 0x44))

add_text(s, "Summary", 0.6, 0.6, 12.0, 0.7,
         size=36, bold=True, color=WHITE)
add_text(s,
    "An end-to-end automated hydrocarbon tank reconciliation pipeline on SAP BTP,\n"
    "eliminating 2–4 hours of daily manual work per terminal.",
    0.6, 1.35, 12.0, 0.9, size=15, color=MID_GREY, italic=True)

pillars = [
    ("IS-OIL\nIntegration", "First live OGS/650 → BTP\nconnection in this landscape"),
    ("Full\nAutomation", "Gauge reading → Material\nDocument in < 2 minutes"),
    ("Governance\n& Audit", "No posting without approval.\nImmutable audit trail."),
    ("Intelligent\nAlerts", "OK/FLAG/URGENT routing to\nright roles, right channels"),
    ("Single\nCockpit", "One React dashboard for all\n5 roles — no spreadsheets"),
]
for i, (title, desc) in enumerate(pillars):
    x = 0.55 + i * 2.44
    add_rect(s, x, 2.5, 2.18, 2.35, fill=RGBColor(0x00, 0x2A, 0x55), line=SAP_BLUE, line_w=0.8)
    add_text(s, title, x+0.1, 2.58, 1.98, 0.72,
             size=12, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    add_text(s, desc, x+0.1, 3.35, 1.98, 0.95,
             size=9.5, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 0.6, 5.35, 12.13, 0.05, fill=SAP_BLUE)
add_text(s,
    "n8n workflow  ·  SAP CAP  ·  React + SAP UI5  ·  IS-OIL HPM  ·  Cloud Connector  ·  XSUAA  ·  BTP ANS",
    0.6, 5.55, 12.0, 0.45, size=11, color=MID_GREY, align=PP_ALIGN.CENTER)
add_text(s, "Terminal Operations / Supply Chain  |  2026",
         0.6, 6.1, 12.0, 0.4, size=11, color=RGBColor(0x66, 0x88, 0xAA),
         align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────
out = "/Users/i023725/ong_tank_reconciliation/Tank_Reconciliation_Solution.pptx"
prs.save(out)
print(f"Saved: {out}")
