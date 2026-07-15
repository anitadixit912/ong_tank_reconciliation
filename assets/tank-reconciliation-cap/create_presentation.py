from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# SAP Colors
SAP_BLUE       = RGBColor(0x00, 0x6D, 0xC1)
SAP_DARK_BLUE  = RGBColor(0x00, 0x3B, 0x6E)
SAP_TEAL       = RGBColor(0x04, 0x7D, 0x9C)
SAP_GREEN      = RGBColor(0x18, 0x84, 0x25)
SAP_ORANGE     = RGBColor(0xE7, 0x6E, 0x00)
SAP_RED        = RGBColor(0xBB, 0x00, 0x00)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY     = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GREY      = RGBColor(0x33, 0x33, 0x33)
MID_GREY       = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h, font_size=12, bold=False, color=DARK_GREY,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.1, fill=SAP_DARK_BLUE)
    add_text_box(slide, title, 0.3, 0.1, 12, 0.6, font_size=24, bold=True,
                 color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.3, 0.65, 12, 0.4, font_size=12,
                     color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.LEFT)

def add_footer(slide, text="Hydrocarbon Tank Stock Reconciliation  |  SAP BTP + IS-OIL OGS/650"):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=SAP_DARK_BLUE)
    add_text_box(slide, text, 0.3, 7.12, 12, 0.3, font_size=9,
                 color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.LEFT)

def add_box(slide, title, body_lines, l, t, w, h,
            title_bg=SAP_BLUE, body_bg=LIGHT_GREY, title_color=WHITE, body_color=DARK_GREY):
    add_rect(slide, l, t, w, 0.38, fill=title_bg)
    add_text_box(slide, title, l+0.08, t+0.04, w-0.16, 0.3,
                 font_size=11, bold=True, color=title_color)
    add_rect(slide, l, t+0.38, w, h-0.38, fill=body_bg, line=SAP_BLUE, line_w=0.5)
    body_text = "\n".join(body_lines)
    add_text_box(slide, body_text, l+0.1, t+0.42, w-0.2, h-0.5,
                 font_size=9.5, color=body_color, wrap=True)

# ─── SLIDE 1: TITLE ───────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
add_rect(s1, 0, 0, 13.33, 7.5, fill=SAP_DARK_BLUE)
add_rect(s1, 0, 4.8, 13.33, 2.7, fill=SAP_BLUE)

add_text_box(s1, "🛢  Hydrocarbon Tank Stock", 1, 1.2, 11, 1.0,
             font_size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text_box(s1, "Reconciliation", 1, 2.1, 11, 0.9,
             font_size=38, bold=True, color=RGBColor(0x7E, 0xC8, 0xE3), align=PP_ALIGN.LEFT)
add_text_box(s1, "End-to-End Automated Pipeline  |  SAP BTP + IS-OIL OGS/650",
             1, 3.1, 11, 0.5, font_size=16, color=RGBColor(0xCC, 0xDD, 0xEE))
add_rect(s1, 1, 3.65, 4, 0.05, fill=SAP_ORANGE)

add_text_box(s1, "From tank gauge reading to S/4HANA material document — fully automated",
             1, 5.1, 11, 0.5, font_size=14, color=WHITE, italic=True)
add_text_box(s1, "SAP BTP  ·  CAP  ·  n8n  ·  React  ·  IS-OIL HPM  ·  Cloud Connector",
             1, 5.7, 11, 0.4, font_size=11, color=RGBColor(0xAA, 0xBB, 0xCC))

# ─── SLIDE 2: BUSINESS PROBLEM ────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
add_header(s2, "The Business Problem", "Why manual reconciliation is broken")
add_footer(s2)

add_text_box(s2, "Every day, terminal operators must answer two questions:",
             0.4, 1.25, 12, 0.35, font_size=12, color=DARK_GREY)
add_rect(s2, 0.4, 1.65, 5.8, 0.55, fill=SAP_BLUE)
add_text_box(s2, "\"How much product is physically in the tank?\"",
             0.5, 1.7, 5.6, 0.45, font_size=13, bold=True, color=WHITE)
add_rect(s2, 6.8, 1.65, 5.8, 0.55, fill=SAP_TEAL)
add_text_box(s2, "\"How much does SAP say should be there?\"",
             6.9, 1.7, 5.6, 0.45, font_size=13, bold=True, color=WHITE)

add_text_box(s2, "The gap = Reconciliation Variance  →  undetected = financial loss, compliance risk, safety risk",
             0.4, 2.35, 12.4, 0.35, font_size=11, italic=True, color=MID_GREY)

problems = [
    ("⏱  Slow", ["2–4 hours of manual work", "per terminal, per day"]),
    ("❌  Error-Prone", ["Wrong VCF table", "Wrong material number", "Calculation mistakes"]),
    ("🔒  Not Auditable", ["Email chains & spreadsheets", "No tamper-evident trail", "No actor/timestamp"]),
    ("🧩  Fragmented", ["Data in ATG, SAP, Excel,", "and email simultaneously"]),
    ("📈  Not Scalable", ["Adding a terminal =", "more manual work"]),
]
for i, (title, lines) in enumerate(problems):
    add_box(s2, title, lines, 0.3 + i*2.54, 2.85, 2.4, 1.9,
            title_bg=SAP_DARK_BLUE if i % 2 == 0 else SAP_BLUE)

add_rect(s2, 0.3, 5.0, 12.4, 0.9, fill=RGBColor(0xFF, 0xF3, 0xE0))
add_rect(s2, 0.3, 5.0, 0.08, 0.9, fill=SAP_ORANGE)
add_text_box(s2, "The hidden complexity — IS-OIL HPM extends SAP with tank-specific data (OIB_TANKDIP, strapping tables, VCF conversion) that has NO public OData APIs. Building a digital solution requires deep IS-OIL expertise.",
             0.5, 5.05, 12, 0.8, font_size=10.5, color=DARK_GREY)

# ─── SLIDE 3: THE SOLUTION ────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
add_header(s3, "The Solution", "Fully automated — from tank gauge to material document")
add_footer(s3)

befores = [
    "Operator reads dip manually, writes on paper",
    "Book stock looked up manually in SAP",
    "VCF correction done with Excel/lookup tables",
    "Delta calculated in spreadsheet",
    "Supervisor approval via email",
    "Goods movement posted manually in MIGO",
    "Report emailed manually",
    "No audit trail",
]
afters = [
    "IS-OIL OIB_TANKDIP read automatically via ZTANK_DIP_SRV_SRV",
    "RELSTOCK fetched directly from IS-OIL dip record",
    "VCF Calculator applies automatically, ASTM fallback if needed",
    "Variance Engine computes delta and classifies in milliseconds",
    "Approval Queue in CAP dashboard with full audit record",
    "Material Document created automatically via API_MATERIAL_DOCUMENT_SRV",
    "PDF generated and sent to Email + MS Teams automatically",
    "Every milestone, decision and posting recorded with timestamp",
]

add_rect(s3, 0.3, 1.15, 5.8, 0.35, fill=SAP_RED)
add_text_box(s3, "BEFORE", 0.3, 1.18, 5.8, 0.3, font_size=11, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s3, 7.1, 1.15, 5.8, 0.35, fill=SAP_GREEN)
add_text_box(s3, "AFTER", 7.1, 1.18, 5.8, 0.3, font_size=11, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

for i, (b, a) in enumerate(zip(befores, afters)):
    bg = LIGHT_GREY if i % 2 == 0 else WHITE
    add_rect(s3, 0.3, 1.5 + i*0.6, 5.8, 0.58, fill=bg, line=RGBColor(0xDD, 0xDD, 0xDD), line_w=0.3)
    add_text_box(s3, "✗  " + b, 0.4, 1.52 + i*0.6, 5.6, 0.55, font_size=9, color=SAP_RED)
    add_rect(s3, 7.1, 1.5 + i*0.6, 5.8, 0.58, fill=bg, line=RGBColor(0xDD, 0xDD, 0xDD), line_w=0.3)
    add_text_box(s3, "✓  " + a, 7.2, 1.52 + i*0.6, 5.6, 0.55, font_size=9, color=SAP_GREEN)
    add_rect(s3, 6.2, 1.55 + i*0.6, 0.7, 0.45, fill=SAP_ORANGE)
    add_text_box(s3, "→", 6.2, 1.55 + i*0.6, 0.7, 0.45, font_size=18, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)

# ─── SLIDE 4: ARCHITECTURE ────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
add_header(s4, "Architecture", "Three layers working together")
add_footer(s4)

# Layer 1 - OGS
add_rect(s4, 0.2, 1.2, 3.5, 5.5, fill=RGBColor(0xE8, 0xF4, 0xFD), line=SAP_BLUE, line_w=1)
add_rect(s4, 0.2, 1.2, 3.5, 0.4, fill=SAP_DARK_BLUE)
add_text_box(s4, "🏭  Layer 1 — OGS/650", 0.3, 1.23, 3.3, 0.35, font_size=10, bold=True, color=WHITE)
add_text_box(s4, "IS-OIL S/4HANA Private Cloud\n(Source of Truth)", 0.3, 1.65, 3.3, 0.5, font_size=9, color=MID_GREY, italic=True)
add_box(s4, "🛢 OIB_TANKDIP", ["Physical dip readings", "Book stock (RELSTOCK)", "SOCNR-based tank ID"], 0.3, 2.25, 3.3, 1.2, title_bg=SAP_TEAL)
add_box(s4, "📡 ZTANK_DIP_SRV_SRV", ["Custom OData Service", "Exposes dip data to BTP", "(built from scratch)"], 0.3, 3.55, 3.3, 1.1, title_bg=SAP_BLUE)
add_box(s4, "🔧 Z_TANK_RECON_TRIGGER_RUN", ["ABAP FM to trigger BTP", "Fetches XSUAA token", "Calls CAP endpoint"], 0.3, 4.75, 3.3, 1.1, title_bg=SAP_BLUE)

# Arrow + Cloud Connector
add_rect(s4, 3.7, 3.3, 1.5, 0.5, fill=SAP_ORANGE)
add_text_box(s4, "☁  Cloud\nConnector", 3.72, 3.32, 1.46, 0.46, font_size=8, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s4, "APAC_DEV10", 3.72, 3.82, 1.46, 0.25, font_size=7, color=MID_GREY, align=PP_ALIGN.CENTER)
add_text_box(s4, "→", 5.1, 3.35, 0.3, 0.4, font_size=20, bold=True, color=SAP_ORANGE, align=PP_ALIGN.CENTER)

# Layer 2 - BTP
add_rect(s4, 5.4, 1.2, 4.5, 5.5, fill=RGBColor(0xF0, 0xF8, 0xF0), line=SAP_GREEN, line_w=1)
add_rect(s4, 5.4, 1.2, 4.5, 0.4, fill=SAP_GREEN)
add_text_box(s4, "⚙️  Layer 2 — SAP BTP", 5.5, 1.23, 4.3, 0.35, font_size=10, bold=True, color=WHITE)
add_box(s4, "🔄 n8n Workflow", ["1. Fetch IS-OIL dip data", "2. VCF correction", "3. Variance calculation", "4. Post Material Document", "5. Send alerts & PDF report"], 5.5, 1.65, 4.2, 1.8, title_bg=SAP_TEAL)
add_box(s4, "📦 CAP Application", ["Stores all run history", "Stores variance records", "Approval state machine", "Full audit trail (M1-M6)", "OData APIs for dashboard"], 5.5, 3.55, 4.2, 1.8, title_bg=SAP_DARK_BLUE)
add_box(s4, "🔐 XSUAA", ["Role-based access control", "M2M token for OGS calls", "User auth for dashboard"], 5.5, 5.45, 4.2, 0.9, title_bg=SAP_BLUE)

# Arrow to Layer 3
add_text_box(s4, "→", 9.85, 3.35, 0.3, 0.4, font_size=20, bold=True, color=SAP_GREEN, align=PP_ALIGN.CENTER)

# Layer 3 - Dashboard
add_rect(s4, 10.15, 1.2, 3.0, 5.5, fill=RGBColor(0xFD, 0xF5, 0xE8), line=SAP_ORANGE, line_w=1)
add_rect(s4, 10.15, 1.2, 3.0, 0.4, fill=SAP_ORANGE)
add_text_box(s4, "👁  Layer 3 — Dashboard", 10.25, 1.23, 2.8, 0.35, font_size=10, bold=True, color=WHITE)
add_box(s4, "📊 Dashboard", ["All roles — run status", "KPI tiles, variance table"], 10.25, 1.65, 2.7, 1.0, title_bg=SAP_BLUE)
add_box(s4, "✅ Approval Queue", ["Supervisor only", "Approve / Reject URGENT"], 10.25, 2.75, 2.7, 0.9, title_bg=SAP_DARK_BLUE)
add_box(s4, "📋 Audit Trail", ["All roles — read only", "M1–M6 milestones"], 10.25, 3.75, 2.7, 0.9, title_bg=SAP_TEAL)
add_box(s4, "⚙️ Configuration", ["Admin only", "Tank thresholds CRUD"], 10.25, 4.75, 2.7, 0.9, title_bg=SAP_BLUE)
add_box(s4, "💬 AI Assistant", ["All roles", "Chat with reconciliation data"], 10.25, 5.75, 2.7, 0.85, title_bg=SAP_TEAL)

# ─── SLIDE 5: 6 MILESTONES ────────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
add_header(s5, "The Flow — 6 Milestones", "From dip reading to report distribution")
add_footer(s5)

milestones = [
    ("M1", "Data Ingestion", SAP_TEAL,
     ["Read OIB_TANKDIP via ZTANK_DIP_SRV_SRV",
      "Physical qty (QUAN_SKU) + Book stock (RELSTOCK)",
      "Missing dip → halt + URGENT alert"]),
    ("M2", "VCF Correction", SAP_BLUE,
     ["Convert gross → net standard volume",
      "Hydrocarbon Qty Conversion API",
      "ASTM D1250 fallback if API unavailable"]),
    ("M3", "Variance Calculation", SAP_DARK_BLUE,
     ["Delta = Physical − Book Stock",
      "OK (≤0.10%)  FLAG (≤0.25%)  URGENT (>0.25%)",
      "Per-tank configurable thresholds"]),
    ("M4", "Approval Decision", SAP_ORANGE,
     ["URGENT tanks held in CAP",
      "Supervisor: Approve or Reject",
      "No posting without approval record"]),
    ("M5", "Goods Movement Posting", SAP_GREEN,
     ["OK/FLAG: auto-post",
      "URGENT: post after approval",
      "551=shrinkage  552=gain  Doc ID written back"]),
    ("M6", "Report Distribution", SAP_TEAL,
     ["Per-tank PDF generated",
      "Email + MS Teams distribution",
      "BTP Alert Notification Service"]),
]

for i, (code, title, color, lines) in enumerate(milestones):
    col = i % 3
    row = i // 3
    lx = 0.25 + col * 4.35
    ty = 1.2 + row * 2.95
    add_rect(s5, lx, ty, 4.1, 0.55, fill=color)
    add_text_box(s5, code, lx + 0.1, ty + 0.05, 0.6, 0.45, font_size=18, bold=True, color=WHITE)
    add_text_box(s5, title, lx + 0.72, ty + 0.1, 3.2, 0.38, font_size=13, bold=True, color=WHITE)
    add_rect(s5, lx, ty + 0.55, 4.1, 2.3, fill=LIGHT_GREY, line=color, line_w=0.8)
    for j, line in enumerate(lines):
        add_text_box(s5, "•  " + line, lx + 0.15, ty + 0.65 + j * 0.62, 3.8, 0.58,
                     font_size=10, color=DARK_GREY)
    if col < 2:
        add_text_box(s5, "→", lx + 4.1, ty + 0.15, 0.25, 0.35, font_size=16,
                     bold=True, color=MID_GREY, align=PP_ALIGN.CENTER)

# ─── SLIDE 6: CHALLENGES ──────────────────────────────────────────────────────
s6 = prs.slides.add_slide(blank_layout)
add_header(s6, "Challenges We Faced", "Real obstacles — and how we solved them")
add_footer(s6)

challenges = [
    ("🔌", "IS-OIL has no public APIs",
     "No OData for OIB_TANKDIP. Built ZTANK_DIP_SRV_SRV from scratch using SEGW + OIIC_DIP_READ_TANKDIPS FM."),
    ("💻", "SAP GUI crashes on Mac",
     "SE37/SE38 threw CNTL_ERROR on macOS. Switched to web-based SAP GUI to write ABAP."),
    ("🌐", "Cloud Connector proxy protocol",
     "ECONNREFUSED → 503 → 501 → 200. Fixed by adding Proxy-Authorization + SAP-Connectivity-SCC-Location_ID headers."),
    ("🔐", "XSUAA tenant mode locked",
     "tenant-mode: shared blocked AppRouter. Fixed via UAA_SERVICE_NAME + TENANT_HOST_PATTERN env vars."),
    ("🔍", "OData filter case sensitivity",
     "'Socnr' vs 'SOCNR' — all tanks returned same wrong record. Fixed ABAP CASE statement to use uppercase."),
    ("🗺", "IS-OIL data model discovery",
     "No docs for plant/location → SOCNR mapping. Traced through OIISOCISL, OIISOCK, and FM parameters to find it."),
    ("🔑", "OA2C_CONFIG unreachable on Mac",
     "OAuth config redirects to browser — not accessible. Built ABAP FM with full client credentials flow instead."),
]

for i, (icon, title, desc) in enumerate(challenges):
    col = i % 2
    row = i // 2
    if i == 6:
        lx = 3.17
    else:
        lx = 0.3 + col * 6.5
    ty = 1.25 + row * 1.45
    add_rect(s6, lx, ty, 6.1, 1.3, fill=LIGHT_GREY, line=SAP_BLUE, line_w=0.5)
    add_rect(s6, lx, ty, 0.08, 1.3, fill=SAP_ORANGE)
    add_text_box(s6, icon + "  " + title, lx + 0.2, ty + 0.05, 5.8, 0.38,
                 font_size=10.5, bold=True, color=SAP_DARK_BLUE)
    add_text_box(s6, desc, lx + 0.2, ty + 0.45, 5.8, 0.75, font_size=9, color=DARK_GREY)

# ─── SLIDE 7: ACHIEVEMENTS ────────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank_layout)
add_header(s7, "Key Achievements", "What we delivered")
add_footer(s7)

achievements = [
    ("✅", "First IS-OIL → BTP Integration", SAP_GREEN,
     "Live OIB_TANKDIP data flowing from OGS/650 to BTP via Cloud Connector. Verified against source records."),
    ("✅", "Custom OData Services Built", SAP_GREEN,
     "ZTANK_DIP_SRV_SRV + ZTANK_PLANT_SRV_SRV built in OGS/650 — IS-OIL data exposed as OData for the first time."),
    ("✅", "OGS → BTP M2M Integration", SAP_TEAL,
     "Z_TANK_RECON_TRIGGER_RUN ABAP FM: fetches XSUAA token, calls BTP CAP — IS-OIL triggers BTP directly."),
    ("✅", "Full Approval Governance", SAP_BLUE,
     "URGENT variances held in CAP state machine. No Material Document posted without supervisor sign-off."),
    ("✅", "Data Accuracy Verified", SAP_DARK_BLUE,
     "Dashboard values match OIB_TANKDIP exactly — book stock, physical qty, delta verified in SE16N."),
    ("✅", "Role-Based Access Control", SAP_BLUE,
     "XSUAA scopes + React ProtectedRoute — Approvals: Supervisor only, Config: Admin only."),
    ("⚡", "2 Minutes vs. 2–4 Hours", SAP_ORANGE,
     "Full reconciliation pipeline — from tank dip to audit trail — completes in under 2 minutes."),
]

for i, (icon, title, color, desc) in enumerate(achievements):
    col = i % 2
    row = i // 2
    if i == 6:
        lx = 3.17
    else:
        lx = 0.3 + col * 6.5
    ty = 1.25 + row * 1.45
    add_rect(s7, lx, ty, 6.1, 1.3, fill=LIGHT_GREY, line=color, line_w=0.8)
    add_rect(s7, lx, ty, 0.08, 1.3, fill=color)
    add_text_box(s7, icon + "  " + title, lx + 0.2, ty + 0.05, 5.8, 0.38,
                 font_size=10.5, bold=True, color=color)
    add_text_box(s7, desc, lx + 0.2, ty + 0.45, 5.8, 0.75, font_size=9, color=DARK_GREY)

# ─── SLIDE 8: CLOSING ─────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(blank_layout)
add_rect(s8, 0, 0, 13.33, 7.5, fill=SAP_DARK_BLUE)
add_rect(s8, 0, 4.5, 13.33, 3.0, fill=SAP_BLUE)

add_text_box(s8, "Architecture in One Sentence", 1, 1.0, 11, 0.5,
             font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC))
add_rect(s8, 1, 1.55, 11, 0.05, fill=SAP_ORANGE)

sentence = (
    "n8n workflow on BTP orchestrates the pipeline  →  "
    "CAP is the backbone & approval state machine  →  "
    "React Dashboard is the single cockpit  →  "
    "IS-OIL OGS/650 is the source of truth  →  "
    "Cloud Connector bridges on-premise to BTP securely"
)
add_text_box(s8, sentence, 1, 1.7, 11, 1.2, font_size=15, bold=True,
             color=WHITE, wrap=True)

boxes = [
    ("n8n", "Orchestration\nEngine"),
    ("CAP", "Backbone &\nState Machine"),
    ("React", "Single\nCockpit"),
    ("IS-OIL", "Source of\nTruth"),
    ("Cloud\nConnector", "Secure\nBridge"),
]
colors = [SAP_TEAL, SAP_BLUE, SAP_ORANGE, SAP_GREEN, SAP_DARK_BLUE]
for i, ((label, sub), color) in enumerate(zip(boxes, colors)):
    lx = 1.0 + i * 2.27
    add_rect(s8, lx, 3.2, 1.9, 0.9, fill=color)
    add_text_box(s8, label, lx, 3.25, 1.9, 0.4, font_size=13, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s8, sub, lx, 3.65, 1.9, 0.4, font_size=8,
                 color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(s8, "→", lx + 1.9, 3.35, 0.37, 0.5, font_size=18,
                     bold=True, color=SAP_ORANGE, align=PP_ALIGN.CENTER)

add_text_box(s8, "Thank You", 1, 5.0, 11, 0.8,
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s8, "Hydrocarbon Tank Stock Reconciliation  |  SAP BTP + IS-OIL OGS/650",
             1, 5.85, 11, 0.4, font_size=12,
             color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

# Save
output = "/Users/i023725/ong_tank_reconciliation/assets/tank-reconciliation-cap/Tank_Reconciliation_Presentation.pptx"
prs.save(output)
print("Saved:", output)
